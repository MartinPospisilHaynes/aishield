"""
AIshield.cz — PPTX Generator
Generování profesionálních PowerPoint prezentací pro AI Act školení.

Strategie: Master šablona (branded design) + python-pptx plní placeholdery.
Claude / AI NEGENERUJE layout — pouze obsah (JSON → slide data).

Pokud šablona aishield_template.pptx neexistuje, vytvoří se programově
s AIshield branding (fuchsia/cyan gradient, tmavé pozadí).
"""

import io
import logging
from datetime import datetime, timezone

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger(__name__)

# ── AIshield Brand Colors ──
COLOR_BG_DARK = RGBColor(0x0F, 0x17, 0x2A)       # #0f172a
COLOR_SURFACE = RGBColor(0x1E, 0x29, 0x3B)        # #1e293b
COLOR_PRIMARY = RGBColor(0xE8, 0x79, 0xF9)        # #e879f9 fuchsia
COLOR_SECONDARY = RGBColor(0x22, 0xD3, 0xEE)      # #22d3ee cyan
COLOR_TEXT = RGBColor(0xF1, 0xF5, 0xF9)            # #f1f5f9
COLOR_MUTED = RGBColor(0x94, 0xA3, 0xB8)           # #94a3b8
COLOR_HIGH = RGBColor(0xEF, 0x44, 0x44)            # #ef4444 red
COLOR_LIMITED = RGBColor(0xF5, 0x9E, 0x0B)         # #f59e0b amber
COLOR_MINIMAL = RGBColor(0x22, 0xC5, 0x5E)         # #22c55e green
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def _set_slide_bg(slide, color=COLOR_BG_DARK):
    """Nastaví pozadí slidu na tmavou barvu."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, left, top, width, height, text, font_size=18,
                  color=COLOR_TEXT, bold=False, alignment=PP_ALIGN.LEFT):
    """Přidá textbox s AIshield stylem."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = alignment
    return txBox


def _add_bullet_list(slide, left, top, width, height, items, font_size=16):
    """Přidá odrážkový seznam."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLOR_TEXT
        p.font.name = "Calibri"
        p.space_after = Pt(8)
        p.level = 0

    return txBox


def _add_branded_header(slide, title_text, client_info=None):
    """Přidá AIshield branded header na slide + klientské info vpravo nahoře."""
    # Brand logo text
    _add_text_box(
        slide,
        left=Inches(0.6), top=Inches(0.3),
        width=Inches(3), height=Inches(0.5),
        text="AIshield.cz",
        font_size=14, color=COLOR_PRIMARY, bold=True,
    )

    # Client info — vpravo nahoře (firma, osoba)
    if client_info:
        info_parts = []
        if client_info.get("company"):
            info_parts.append(client_info["company"])
        if client_info.get("person"):
            info_parts.append(client_info["person"])
        if info_parts:
            _add_text_box(
                slide,
                left=Inches(8), top=Inches(0.2),
                width=Inches(5), height=Inches(0.6),
                text="  •  ".join(info_parts),
                font_size=11, color=COLOR_MUTED,
                alignment=PP_ALIGN.RIGHT,
            )

    # Fuchsia accent line
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(0.85),
        Inches(2), Inches(0.04),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.fill.background()

    # Slide title
    _add_text_box(
        slide,
        left=Inches(0.6), top=Inches(1.1),
        width=Inches(11), height=Inches(0.8),
        text=title_text,
        font_size=28, color=COLOR_TEXT, bold=True,
    )


def _add_client_footer(slide, client_info):
    """Přidá zápatí s kontaktními údaji klienta na slide."""
    if not client_info:
        return

    parts = []
    if client_info.get("company"):
        parts.append(client_info["company"])
    if client_info.get("person"):
        parts.append(client_info["person"])
    if client_info.get("phone"):
        parts.append(client_info["phone"])
    if client_info.get("email"):
        parts.append(client_info["email"])

    if parts:
        footer_text = "  •  ".join(parts)
        _add_text_box(
            slide,
            left=Inches(0.5), top=Inches(6.9),
            width=Inches(12), height=Inches(0.4),
            text=footer_text,
            font_size=9, color=COLOR_MUTED,
            alignment=PP_ALIGN.CENTER,
        )


def _create_title_slide(prs, company_name, subtitle="", client_info=None):
    """Titulní slide s velkým logem a názvem firmy."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    _set_slide_bg(slide)

    # Centered brand
    _add_text_box(
        slide,
        left=Inches(2), top=Inches(1.5),
        width=Inches(9), height=Inches(1),
        text="AIshield.cz",
        font_size=48, color=COLOR_PRIMARY, bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # Accent line
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5), Inches(2.7),
        Inches(3), Inches(0.05),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_SECONDARY
    shape.line.fill.background()

    # Title
    _add_text_box(
        slide,
        left=Inches(1.5), top=Inches(3),
        width=Inches(10), height=Inches(1),
        text=f"Školení AI Literacy — {company_name}",
        font_size=32, color=COLOR_TEXT, bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # Subtitle
    sub = subtitle or "Povinné školení dle čl. 4 Nařízení (EU) 2024/1689 (AI Act)"
    _add_text_box(
        slide,
        left=Inches(2), top=Inches(4.2),
        width=Inches(9), height=Inches(0.6),
        text=sub,
        font_size=16, color=COLOR_MUTED,
        alignment=PP_ALIGN.CENTER,
    )

    # Footer info
    now = datetime.now(timezone.utc).strftime("%d. %m. %Y")
    _add_text_box(
        slide,
        left=Inches(2), top=Inches(5.5),
        width=Inches(9), height=Inches(0.5),
        text=f"Vygenerováno: {now}  •  Rozsah: 2–3 hodiny  •  Cílová skupina: Všichni zaměstnanci",
        font_size=12, color=COLOR_MUTED,
        alignment=PP_ALIGN.CENTER,
    )

    # Client footer
    _add_client_footer(slide, client_info)


def _create_content_slide(prs, title, bullets, client_info=None):
    """Standardní obsahový slide s nadpisem a odrážkami."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _set_slide_bg(slide)
    _add_branded_header(slide, title, client_info=client_info)
    _add_bullet_list(
        slide,
        left=Inches(0.8), top=Inches(2.2),
        width=Inches(11), height=Inches(4.5),
        items=bullets,
        font_size=18,
    )
    _add_client_footer(slide, client_info)


def _create_risk_slide(prs, findings, client_info=None):
    """Slide s přehledem rizik — barevně kódované."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_branded_header(slide, "AI systémy v naší firmě — přehled rizik", client_info=client_info)

    if not findings:
        _add_text_box(
            slide,
            left=Inches(0.8), top=Inches(2.5),
            width=Inches(11), height=Inches(1),
            text="Žádné AI systémy nebyly detekovány na webu firmy.",
            font_size=20, color=COLOR_MUTED,
        )
        return

    y_pos = Inches(2.2)
    for i, f in enumerate(findings[:8]):  # Max 8 per slide
        name = f.get("name", "AI systém")
        risk = f.get("risk_level", "minimal")
        article = f.get("ai_act_article", "")

        color_map = {"high": COLOR_HIGH, "limited": COLOR_LIMITED, "minimal": COLOR_MINIMAL}
        risk_color = color_map.get(risk, COLOR_MUTED)

        # Risk dot
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.8), y_pos + Pt(4),
            Inches(0.15), Inches(0.15),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = risk_color
        dot.line.fill.background()

        # Finding text
        risk_cs = {"high": "VYSOKÉ", "limited": "OMEZENÉ", "minimal": "MINIMÁLNÍ", "none": "Mimo klasifikaci"}.get(risk, risk)
        _add_text_box(
            slide,
            left=Inches(1.1), top=y_pos,
            width=Inches(11), height=Inches(0.4),
            text=f"{name}  —  {risk_cs} riziko  •  {article}",
            font_size=16, color=COLOR_TEXT,
        )

        y_pos += Inches(0.55)

    _add_client_footer(slide, client_info)


def _create_disclaimer_slide(prs):
    """Závěrečný slide s disclaimerem."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    _add_text_box(
        slide,
        left=Inches(2), top=Inches(1.5),
        width=Inches(9), height=Inches(1),
        text="AIshield.cz",
        font_size=36, color=COLOR_PRIMARY, bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    _add_text_box(
        slide,
        left=Inches(1.5), top=Inches(3),
        width=Inches(10), height=Inches(2.5),
        text=(
            "Dokumenty vygenerované platformou AIshield.cz jsou vytvořeny na základě "
            "automatizované analýzy a informací poskytnutých klientem.\n\n"
            "AIshield.cz poskytuje compliance dokumentaci jako technickou pomůcku "
            "pro splnění požadavků Nařízení (EU) 2024/1689 (AI Act), "
            "nikoliv jako právní poradenství ve smyslu zák. č. 85/1996 Sb.\n\n"
            "info@aishield.cz  •  +420 732 716 141  •  Martin Haynes, IČO: 17889251"
        ),
        font_size=14, color=COLOR_MUTED,
        alignment=PP_ALIGN.CENTER,
    )


# ══════════════════════════════════════════════════════════════════════
# HLAVNÍ FUNKCE
# ══════════════════════════════════════════════════════════════════════

def generate_training_pptx(data: dict) -> bytes:
    """
    Generuje kompletní PPTX prezentaci pro AI Act školení.

    Args:
        data: dict s company_name, findings, questionnaire data, oversight atd.

    Returns:
        bytes — obsah PPTX souboru
    """
    company = data.get("company_name", "Firma")
    findings = data.get("findings", [])
    ai_declared = data.get("ai_systems_declared", [])
    oversight = data.get("oversight_person", {})
    training = data.get("training", {})
    incident = data.get("incident", {})
    data_protection = data.get("data_protection", {})
    prohibited = data.get("prohibited_systems", {})
    industry = data.get("q_company_industry", "")
    company_size = data.get("q_company_size", "")

    # ── Client info pro záhlaví/zápatí na každém slidu ──
    client_info = {"company": company}
    # Kontaktní osoba — z oversight nebo z dotazníku
    if oversight.get("name"):
        client_info["person"] = oversight["name"]
    if oversight.get("phone"):
        client_info["phone"] = oversight["phone"]
    if oversight.get("email"):
        client_info["email"] = oversight["email"]
    elif data.get("contact_email"):
        client_info["email"] = data["contact_email"]
    elif data.get("q_company_contact_email"):
        client_info["email"] = data["q_company_contact_email"]

    # Celkový počet AI nástrojů (web scan + dotazník)
    all_tool_names = set()
    for f in findings:
        name = f.get("name")
        if name:
            all_tool_names.add(name)
    for s in ai_declared:
        name = s.get("tool_name") or s.get("key") or "AI systém"
        if name:
            all_tool_names.add(name)
    all_tool_names.discard(None)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ── Slide 1: Titulní ──
    subtitle_parts = ["Povinné školení dle čl. 4 Nařízení (EU) 2024/1689 (AI Act)"]
    if industry:
        subtitle_parts.append(f"Obor: {industry}")
    if company_size:
        subtitle_parts.append(f"Velikost firmy: {company_size}")
    _create_title_slide(prs, company, subtitle="  •  ".join(subtitle_parts), client_info=client_info)

    # ── Slide 2: Agenda ──
    _create_content_slide(prs, "Agenda školení", [
        "Modul 1 — Co je umělá inteligence (30 min)",
        "Modul 2 — EU AI Act v kostce (45 min)",
        "Modul 3 — AI v naší firmě (30 min)",
        "Modul 4 — Bezpečné používání AI v praxi (30 min)",
        "Modul 5 — Naše povinnosti (15 min)",
        "Modul 6 — Test a certifikace (15 min)",
    ], client_info=client_info)

    # ── Slide 3: Co je AI ──
    _create_content_slide(prs, "Modul 1 — Co je umělá inteligence", [
        "Definice AI — co to je a co to není",
        "Typy AI: generativní AI, prediktivní modely, expertní systémy",
        "Příklady AI v každodenním životě (navigace, doporučení, chatboty)",
        "AI vs. automatizace — jaký je rozdíl?",
        "Demonstrace: live ukázka ChatGPT / Claude",
    ], client_info=client_info)

    # ── Slide 4: AI Act overview ──
    _create_content_slide(prs, "Modul 2 — EU AI Act v kostce", [
        "Proč EU reguluje AI — ochrana základních práv",
        "Nařízení (EU) 2024/1689 — plná účinnost 2. 8. 2026",
        "4 kategorie rizik: nepřijatelné → vysoké → omezené → minimální",
        "Zakázané praktiky (čl. 5) — co NESMÍME dělat",
        "Povinnosti transparentnosti (čl. 50)",
        "Povinnost AI gramotnosti (čl. 4) — proto jsme tady",
    ], client_info=client_info)

    # ── Slide 5: Rizikové kategorie ──
    _create_content_slide(prs, "4 kategorie rizik AI Act", [
        "🔴 NEPŘIJATELNÉ — zakázáno (social scoring, manipulace, biometrie v reálném čase)",
        "🟠 VYSOKÉ RIZIKO — přísná regulace (HR nábor, credit scoring, zdravotnictví, justice)",
        "🟡 OMEZENÉ RIZIKO — transparentnost (chatboty, deepfakes, generovaný obsah)",
        "🟢 MINIMÁLNÍ RIZIKO — bez povinností (spam filtry, doporučovací algoritmy)",
        "",
        "Pokuty: až 35 mil. EUR nebo 7 % celosvětového obratu firmy",
    ], client_info=client_info)

    # ── Slide 6: Zakázané praktiky ──
    _create_content_slide(prs, "Zakázané AI praktiky (čl. 5)", [
        "Social scoring — hodnocení lidí na základě chování",
        "Manipulativní AI — podprahové ovlivňování rozhodnutí",
        "Biometrická identifikace v reálném čase na veřejnosti",
        "Prediktivní policing na základě profilingu",
        "Emotion recognition na pracovišti / ve školství (výjimky: bezpečnost)",
        "Vytváření databází obličejů ze scraping internetu",
    ], client_info=client_info)

    # ── Slide 7: AI v naší firmě — personalizovaný ──
    slide7_bullets = []

    if all_tool_names:
        tools_str = ", ".join(sorted(all_tool_names)[:6])
        if len(all_tool_names) > 6:
            tools_str += f" a dalších {len(all_tool_names) - 6}"
        slide7_bullets.append(f"Používáme {len(all_tool_names)} AI nástrojů: {tools_str}")
    else:
        slide7_bullets.append(f"Detekovány AI systémy na webu firmy ({len(findings)} nálezů)")

    if oversight.get("has_person") and oversight.get("name"):
        slide7_bullets.append(f"Odpovědná osoba za AI: {oversight['name']} ({oversight.get('role', 'AI Officer')})")
    else:
        slide7_bullets.append("Odpovědná osoba za AI — přiřadit co nejdříve (povinnost dle čl. 14)")

    slide7_bullets.extend([
        "Registr AI systémů — interní dokument (součást Compliance Kitu)",
        "Interní AI politika — co smíme a co ne",
        "Pravidla pro vkládání dat do AI nástrojů",
    ])

    _create_content_slide(prs, f"Modul 3 — AI v {company}", slide7_bullets, client_info=client_info)

    # ── Slide 8: Přehled rizik z dat ──
    # Kombinovat web scan findings + declared systems
    from backend.documents.unified_pdf import QUESTIONNAIRE_RISK_MAP
    combined_findings = list(findings)
    existing_names = {(f.get("name") or "").lower() for f in findings}
    for sys in ai_declared:
        tool_name = sys.get("tool_name") or sys.get("key") or "AI systém"
        if tool_name.lower() not in existing_names:
            # Použít QUESTIONNAIRE_RISK_MAP místo hardcoded "limited"
            sys_key = sys.get("key", "")
            risk_level = QUESTIONNAIRE_RISK_MAP.get(sys_key, "limited")
            article_map = {
                "high": "Příloha III — vysoce rizikový systém",
                "limited": "čl. 50 — transparentnost",
                "minimal": "dobrovolné best practices",
            }
            combined_findings.append({
                "name": tool_name,
                "risk_level": risk_level,
                "ai_act_article": article_map.get(risk_level, "čl. 50 — transparentnost"),
            })
    _create_risk_slide(prs, combined_findings, client_info=client_info)

    # ── Slide 9: Bezpečné používání AI ──
    safe_bullets = [
        "ChatGPT / Claude / Copilot — jak správně a bezpečně používat",
        "Co do AI NIKDY nevkládat (osobní údaje, hesla, smlouvy, interní data)",
        'Ověřování výstupů AI — „trust but verify" (AI halucinuje)',
        "Označování AI-generovaného obsahu pro transparentnost",
    ]
    if incident.get("has_plan"):
        safe_bullets.append("Hlášení incidentů — postupujte dle interního Incident Response Plánu")
    elif oversight.get("has_person") and oversight.get("email"):
        safe_bullets.append(f"Hlášení incidentů → kontaktujte {oversight['email']}")
    else:
        safe_bullets.append("Hlášení incidentů — komu a jak postupovat (stanovit odpovědnou osobu)")
    _create_content_slide(prs, "Modul 4 — Bezpečné používání AI", safe_bullets, client_info=client_info)

    # ── Slide 10: GDPR a AI ──
    gdpr_bullets = [
        "AI zpracovává osobní údaje — platí GDPR",
        "Právní základ: souhlas, oprávněný zájem, plnění smlouvy",
        "DPIA (posouzení vlivu) — povinné pro AI s vysokým rizikem",
        "Právo na vysvětlení automatizovaného rozhodnutí (čl. 22 GDPR)",
        "Minimalizace dat — do AI jen to, co je nezbytné",
    ]
    if data_protection.get("processes_personal_data"):
        gdpr_bullets.append("⚠ Naše AI systémy zpracovávají osobní údaje — zvýšená pozornost!")
    if not data_protection.get("data_in_eu"):
        gdpr_bullets.append("⚠ Některá data mohou být uložena mimo EU — ověřte transfer mechanismy")
    _create_content_slide(prs, "AI a ochrana osobních údajů (GDPR)", gdpr_bullets, client_info=client_info)

    # ── Slide 11: Naše konkrétní povinnosti ──
    duty_bullets = []

    # Zakázané praktiky check
    any_prohibited = any(prohibited.get(k) for k in prohibited) if prohibited else False
    if any_prohibited:
        duty_bullets.append("🔴 POZOR: Identifikovány potenciálně zakázané AI praktiky — vyžadují okamžitou nápravu!")
    else:
        duty_bullets.append("✅ Žádné zakázané AI praktiky nebyly identifikovány")

    # Incident plan
    if incident.get("has_plan"):
        duty_bullets.append("✅ Incident Response Plán existuje")
    else:
        duty_bullets.append("⚠ Incident Response Plán chybí — nutno vytvořit (čl. 73)")

    # AI registr
    if data.get("human_oversight", {}).get("has_register"):
        duty_bullets.append("✅ Registr AI systémů je veden")
    else:
        duty_bullets.append("⚠ Registr AI systémů chybí — nutno zavést (čl. 49)")

    # Data protection
    if data_protection.get("has_vendor_contracts"):
        duty_bullets.append("✅ Smlouvy s dodavateli AI jsou uzavřeny")
    else:
        duty_bullets.append("⚠ Smlouvy s AI dodavateli chybí — ošetřit zpracovatelské smlouvy")

    # Training
    if training.get("has_training"):
        duty_bullets.append("✅ Školení AI gramotnosti probíhá")
    else:
        duty_bullets.append("⚠ Školení AI gramotnosti dosud neproběhlo — právě ho provádíme")

    _create_content_slide(prs, "Modul 5 — Naše konkrétní povinnosti", duty_bullets, client_info=client_info)

    # ── Slide 12: Test ──
    test_bullets = [
        "Krátký test (10 otázek) — ověření porozumění",
        "Minimální skóre pro úspěšné absolvování: 70 %",
        "Certifikát o absolvování školení (evidenční účely)",
        "Uchovávejte evidenci minimálně 5 let (doporučení)",
    ]
    audience = training.get("audience_size", "")
    level = training.get("audience_level", "")
    if audience or level:
        info_parts = []
        if audience:
            info_parts.append(f"počet účastníků: {audience}")
        if level:
            level_labels = {"beginner": "začátečníci", "intermediate": "pokročilí", "advanced": "expertní"}
            info_parts.append(f"úroveň: {level_labels.get(level, level)}")
        test_bullets.append(f"Cílová skupina: {', '.join(info_parts)}")
    else:
        test_bullets.append("")
    test_bullets.append("Tip: Opakujte školení 1× ročně jako refresher")
    _create_content_slide(prs, "Modul 6 — Test a certifikace", test_bullets, client_info=client_info)

    # ── Slide 13: Disclaimer ──
    _create_disclaimer_slide(prs)

    # Export to bytes
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    pptx_bytes = buffer.getvalue()

    logger.info(f"PPTX vygenerováno: {len(pptx_bytes)} bytes, {len(prs.slides)} slidů")
    return pptx_bytes


# ══════════════════════════════════════════════════════════════════════
# HTML → PPTX CONVERTER — LLM-generated slide content to PPTX
# ══════════════════════════════════════════════════════════════════════

import re
from html.parser import HTMLParser


class _SlideExtractor(HTMLParser):
    """
    Parses LLM-generated HTML where each <h2> is a slide title
    and <li>/<p> tags under it are the slide content (bullets/text).
    """

    def __init__(self):
        super().__init__()
        self.slides: list[dict] = []
        self.presentation_title: str = ""
        self._current_slide: dict | None = None
        self._current_tag: str = ""
        self._in_h1 = False
        self._in_h2 = False
        self._in_li = False
        self._in_p = False
        self._text_buf: list[str] = []

    def handle_starttag(self, tag, attrs):
        tag = tag.lower()
        self._current_tag = tag
        if tag == "h1":
            self._in_h1 = True
            self._text_buf = []
        elif tag == "h2":
            # New slide
            if self._current_slide:
                self.slides.append(self._current_slide)
            self._current_slide = {"title": "", "bullets": [], "paragraphs": []}
            self._in_h2 = True
            self._text_buf = []
        elif tag == "li":
            self._in_li = True
            self._text_buf = []
        elif tag == "p" and self._current_slide is not None:
            self._in_p = True
            self._text_buf = []

    def handle_endtag(self, tag):
        tag = tag.lower()
        if tag == "h1" and self._in_h1:
            self._in_h1 = False
            self.presentation_title = "".join(self._text_buf).strip()
        elif tag == "h2" and self._in_h2:
            self._in_h2 = False
            if self._current_slide:
                self._current_slide["title"] = "".join(self._text_buf).strip()
        elif tag == "li" and self._in_li:
            self._in_li = False
            text = "".join(self._text_buf).strip()
            if text and self._current_slide:
                self._current_slide["bullets"].append(text)
        elif tag == "p" and self._in_p:
            self._in_p = False
            text = "".join(self._text_buf).strip()
            if text and self._current_slide:
                self._current_slide["paragraphs"].append(text)
        self._current_tag = ""

    def handle_data(self, data):
        if self._in_h1 or self._in_h2 or self._in_li or self._in_p:
            self._text_buf.append(data)

    def finalize(self):
        if self._current_slide:
            self.slides.append(self._current_slide)
            self._current_slide = None


def html_slides_to_pptx(html_content: str, data: dict) -> bytes:
    """
    Converts LLM-generated HTML slide content to PPTX.

    The HTML is expected to have:
    - <h1> for presentation title
    - <h2> for each slide title
    - <ul><li> for bullet points
    - <p> for descriptive text

    Args:
        html_content: LLM-generated HTML with slide structure
        data: template_data dict for client info

    Returns:
        bytes — PPTX file content
    """
    company = data.get("company_name", "Firma")

    # ── Parse HTML ──
    parser = _SlideExtractor()
    parser.feed(html_content)
    parser.finalize()

    if not parser.slides:
        logger.warning("[PPTX Converter] No <h2> slides found — falling back to template PPTX")
        return generate_training_pptx(data)

    logger.info(f"[PPTX Converter] Extracted {len(parser.slides)} slides from LLM HTML")

    # ── Client info ──
    client_info = {"company": company}
    oversight = data.get("oversight_person", {})
    if oversight.get("name"):
        client_info["person"] = oversight["name"]
    if oversight.get("email"):
        client_info["email"] = oversight["email"]
    elif data.get("contact_email"):
        client_info["email"] = data["contact_email"]
    elif data.get("q_company_contact_email"):
        client_info["email"] = data["q_company_contact_email"]

    # ── Build PPTX ──
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Slide 1: Title slide
    pres_title = parser.presentation_title or f"Školení AI Literacy — {company}"
    first_slide = parser.slides[0] if parser.slides else None

    # If first slide looks like a title slide (fewer bullets), use it
    subtitle_parts = []
    start_idx = 0
    if first_slide and len(first_slide["bullets"]) <= 2:
        # Title slide — use paragraphs and bullets as subtitle
        subtitle_parts = first_slide["paragraphs"] + first_slide["bullets"]
        start_idx = 1

    _create_title_slide(
        prs, company,
        subtitle="  •  ".join(subtitle_parts) if subtitle_parts
                 else "Povinné školení dle čl. 4 Nařízení (EU) 2024/1689 (AI Act)",
        client_info=client_info,
    )

    # Content slides
    for slide_data in parser.slides[start_idx:]:
        title = slide_data["title"]
        bullets = slide_data["bullets"]
        paragraphs = slide_data["paragraphs"]

        # Clean slide number prefix from title if present (e.g. "Slide 5: ...")
        title = re.sub(r'^Slide\s*\d+\s*[:—–-]\s*', '', title, flags=re.IGNORECASE)

        # If it's the last slide and looks like a disclaimer, use disclaimer style
        if slide_data == parser.slides[-1] and (
            "doložka" in title.lower() or "disclaimer" in title.lower()
            or "aishield" in title.lower()
        ):
            _create_disclaimer_slide(prs)
            continue

        # Use bullets if available, otherwise format paragraphs as bullets
        content = bullets if bullets else paragraphs

        if content:
            _create_content_slide(prs, title, content, client_info=client_info)
        else:
            # Empty slide — just title
            _create_content_slide(prs, title, [""], client_info=client_info)

    # Export to bytes
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    pptx_bytes = buffer.getvalue()

    logger.info(f"[PPTX Converter] LLM PPTX vygenerováno: {len(pptx_bytes)} bytes, "
                f"{len(prs.slides)} slidů")
    return pptx_bytes
